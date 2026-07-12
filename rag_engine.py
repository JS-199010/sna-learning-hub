import os
import json
import sqlite3
import numpy as np
import google.generativeai as genai
from pypdf import PdfReader
from pptx import Presentation
import jieba
from rank_bm25 import BM25Okapi
import database

def extract_text(filepath):
    """Extracts raw text from PDF, PPTX, or TXT files."""
    text = ""
    filename = os.path.basename(filepath)
    ext = os.path.splitext(filename)[1].lower()
    
    if ext == '.pdf':
        try:
            reader = PdfReader(filepath)
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text += f"[Page {i+1}]\n{page_text}\n\n"
        except Exception as e:
            raise ValueError(f"Error parsing PDF: {e}")
            
    elif ext == '.pptx':
        try:
            prs = Presentation(filepath)
            for i, slide in enumerate(prs.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                if slide_text.strip():
                    text += f"[Slide {i+1}]\n{slide_text}\n\n"
        except Exception as e:
            raise ValueError(f"Error parsing PPTX: {e}")
            
    elif ext in ['.txt', '.md']:
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        except Exception as e:
            raise ValueError(f"Error parsing text file: {e}")
            
    else:
        raise ValueError(f"Unsupported file format: {ext}")
        
    return text

def get_chunks(text, chunk_size=1000, overlap=200):
    """Chunks the extracted text with a sliding window."""
    if not text.strip():
        return []
        
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start += (chunk_size - overlap)
        
    # If the last chunk is too small and we have other chunks, we can skip or keep it
    if len(chunks) > 1 and len(chunks[-1]) < 100:
        chunks.pop()
        
    return chunks

def _embed_content(content, task_type, api_key):
    """Wrapper that tries multiple Google embedding models in sequence for robustness."""
    genai.configure(api_key=api_key.strip())
    models_to_try = [
        "models/text-embedding-004",
        "models/embedding-001",
        "models/gemini-embedding-001",
        "text-embedding-004",
        "embedding-001"
    ]
    
    last_error = None
    errors = []
    for model_name in models_to_try:
        try:
            res = genai.embed_content(
                model=model_name,
                content=content,
                task_type=task_type
            )
            return res['embedding']
        except Exception as e:
            errors.append(f"{model_name}: {e}")
            last_error = e
            continue
            
    # If all default models fail, let's list the available models to help diagnose
    available_models = []
    try:
        for m in genai.list_models():
            if 'embedContent' in m.supported_generation_methods:
                available_models.append(m.name)
    except Exception as list_err:
        available_models.append(f"(Could not list models: {list_err})")
        
    error_msg = (
        f"Embedding API Error (tried all models).\n\n"
        f"Detailed errors:\n" + "\n".join(errors) + "\n\n"
        f"Your API key supports the following embedding models: {available_models}"
    )
    raise RuntimeError(error_msg)

def get_embeddings_batch(chunks, api_key):
    """Fetch embeddings in batches using the robust embedding helper."""
    embeddings = []
    batch_size = 100
    
    for i in range(0, len(chunks), batch_size):
        batch = chunks[i:i+batch_size]
        try:
            res_embeddings = _embed_content(batch, "retrieval_document", api_key)
            embeddings.extend(res_embeddings)
        except Exception as e:
            raise RuntimeError(f"Failed to fetch embeddings: {e}")
            
    return embeddings

def cosine_similarity(v1, v2):
    """Computes cosine similarity between two numpy vectors."""
    dot_product = np.dot(v1, v2)
    norm_v1 = np.linalg.norm(v1)
    norm_v2 = np.linalg.norm(v2)
    if norm_v1 == 0 or norm_v2 == 0:
        return 0.0
    return float(dot_product / (norm_v1 * norm_v2))

def search_related_chunks(query, filename=None, api_key=None, top_k=5):
    """
    Search database for chunks similar to the query.
    If filename is provided, restricts search to that file; otherwise searches all.
    Uses Hybrid Search (Semantic + BM25).
    """
    if not api_key:
        raise ValueError("API Key is required for similarity search.")
        
    # 1. Embed query using the robust helper
    try:
        query_emb_list = _embed_content(query, "retrieval_query", api_key)
        query_emb = np.array(query_emb_list)
    except Exception as e:
        raise RuntimeError(f"Error embedding search query: {e}")
        
    # 2. Get chunks data using pgvector (if available)
    if filename and filename != "All Files" and filename != "All Uploaded Textbooks":
        target_filename = filename
    else:
        target_filename = None
        
    pgvector_results = database.search_chunks_vector(query_emb_list, match_count=50, filename=target_filename)
    
    semantic_scores = []
    valid_chunks = []
    
    if pgvector_results:
        # RPC succeeded and returned results
        for file_src, chunk_text, sim_score in pgvector_results:
            valid_chunks.append((file_src, chunk_text))
            semantic_scores.append(float(sim_score))
    else:
        # Fallback to local calculation if RPC fails or no vector data (old data)
        if target_filename:
            raw_chunks = database.get_chunks_for_file(target_filename)
            chunks_data = [(target_filename, c[0], c[1]) for c in raw_chunks]
        else:
            chunks_data = database.get_all_chunks()
            
        if not chunks_data:
            return []
            
        for file_src, chunk_text, emb_str in chunks_data:
            try:
                chunk_emb = np.array(json.loads(emb_str))
                s_score = cosine_similarity(query_emb, chunk_emb)
                semantic_scores.append(s_score)
                valid_chunks.append((file_src, chunk_text))
            except Exception:
                continue
                
    if not valid_chunks:
        return []
        
    # 4. Calculate BM25 scores
    # Tokenize texts
    tokenized_query = list(jieba.cut(query))
    tokenized_corpus = [list(jieba.cut(chunk_text)) for _, chunk_text in valid_chunks]
    
    bm25 = BM25Okapi(tokenized_corpus)
    bm25_scores = bm25.get_scores(tokenized_query)
    
    # 5. Normalize and combine (Hybrid)
    def min_max_normalize(scores):
        if len(scores) == 0:
            return []
        min_val = min(scores)
        max_val = max(scores)
        if max_val == min_val:
            return [0.0] * len(scores)
        return [(s - min_val) / (max_val - min_val) for s in scores]
        
    norm_semantic = min_max_normalize(semantic_scores)
    norm_bm25 = min_max_normalize(bm25_scores)
    
    alpha = 0.5  # Semantic weight
    
    scored_chunks = []
    for i in range(len(valid_chunks)):
        hybrid_score = (alpha * norm_semantic[i]) + ((1 - alpha) * norm_bm25[i])
        scored_chunks.append({
            "filename": valid_chunks[i][0],
            "text": valid_chunks[i][1],
            "score": float(hybrid_score),
            "semantic_score": float(norm_semantic[i]),
            "bm25_score": float(norm_bm25[i])
        })
            
    # Sort and pick top K
    scored_chunks.sort(key=lambda x: x["score"], reverse=True)
    return scored_chunks[:top_k]

def clean_json_response(text):
    """Extracts raw JSON content out of markdown code blocks."""
    text = text.strip()
    if "```json" in text:
        text = text.split("```json")[1].split("```")[0]
    elif "```" in text:
        text = text.split("```")[1].split("```")[0]
    return text.strip()

def _generate_with_fallback(api_key, prompt=None, generation_config=None, system_instruction=None, messages=None):
    """Wrapper that tries multiple Google generative models in sequence and executes generate_content."""
    genai.configure(api_key=api_key.strip())
    
    # Priority list of generative models
    models_to_try = [
        "gemini-1.5-flash",
        "gemini-2.5-flash",
        "gemini-2.0-flash",
        "gemini-pro"
    ]
    
    last_err = None
    for model_name in models_to_try:
        try:
            model = genai.GenerativeModel(
                model_name=model_name, 
                system_instruction=system_instruction
            )
            
            # Execute request
            if messages:
                if generation_config:
                    res = model.generate_content(messages, generation_config=generation_config)
                else:
                    res = model.generate_content(messages)
            else:
                if generation_config:
                    res = model.generate_content(prompt, generation_config=generation_config)
                else:
                    res = model.generate_content(prompt)
            return res
        except Exception as e:
            last_err = e
            continue
            
    raise RuntimeError(f"All generative models failed. Last error: {last_err}")

def generate_questions(context_text, topic_query, api_key, num_questions=5, difficulty="中階 (SNA)"):
    """Generates structured questions from the given context using Gemini."""
    import database
    rules = database.load_ai_rules()
    rules_text = "\n".join([f"- {r}" for r in rules]) if rules else "無"
    
    prompt = f"""
You are a highly experienced clinical anesthesia professor. 
Based on the provided textbooks excerpts below, generate {num_questions} multiple-choice questions (MCQs) in Traditional Chinese (繁體中文) related to the topic "{topic_query}".

Target Audience Difficulty: {difficulty}
(Adjust the complexity of vocabulary, clinical depth, and distractor plausibility based on this difficulty level).

【來自總醫師的絕對指令與糾錯 (Teacher's Correction Rules)】:
{rules_text}
(You MUST strictly follow these rules when generating questions and explanations)

Each question must include:
1. "question": The question text, detailing a clinical scenario or concept.
2. "options": An array of exactly 4 option strings (e.g. A, B, C, D details).
3. "answer": The exact string of the correct option (must match one of the items in the "options" array exactly).
4. "explanation": A brief, educational explanation (繁體中文條列式解析) of why the answer is correct and why other options are incorrect, citing terms from the context. Keep it concise.

Provide your response in raw JSON array format matching this schema:
[
  {{
    "question": "...",
    "options": ["Option A", "Option B", "Option C", "Option D"],
    "answer": "Option A",
    "explanation": "..."
  }}
]

Important:
- Return ONLY the JSON array.
- Make the questions clinical, rigorous, and relevant to Student Nurse Anesthetists (SNAs).
- Do not make up facts; rely ONLY on the provided textbook excerpts.

Textbook Excerpts:
{context_text}
"""
    # Use fallback generation helper
    try:
        response = _generate_with_fallback(
            api_key=api_key,
            prompt=prompt,
            generation_config={"response_mime_type": "application/json"}
        )
        json_text = clean_json_response(response.text)
        return json.loads(json_text)
    except Exception as e:
        # Fallback to standard request without response_mime_type if it errors
        try:
            response = _generate_with_fallback(
                api_key=api_key,
                prompt=prompt
            )
            json_text = clean_json_response(response.text)
            return json.loads(json_text)
        except Exception as e2:
            raise RuntimeError(f"Failed to generate questions: {e2}")

def ask_ai_tutor(chat_history, context_text, query, api_key):
    """
    Simulates an AI Study Tutor conversation.
    Uses context to answer user queries with reference.
    """
    import database
    rules = database.load_ai_rules()
    rules_text = "\n".join([f"- {r}" for r in rules]) if rules else "無"
    
    system_instruction = (
        "You are a highly efficient SNA Clinical Anesthesia Tutor. "
        "Explain complex anesthesia concepts concisely. You MUST respond in a brief, highly-structured, and bulleted format (繁體中文條列式). "
        "Avoid long introductory or concluding paragraphs. Keep answers straight to the point and extremely concise. "
        "Use the provided textbook excerpts (Context) as the primary source of truth. If the answer is not in the context, "
        "use your expert clinical knowledge but clearly state it is general clinical advice.\n\n"
        "【來自總醫師的絕對指令與糾錯 (Teacher's Correction Rules)】:\n"
        f"{rules_text}\n"
        "(You MUST strictly follow these rules when responding)"
    )
    
    # Construct prompt with context
    context_prefix = f"【參考教材內容】\n{context_text}\n\n" if context_text else "【無特定教材參考，使用一般臨床麻醉學知識回答】\n\n"
    
    # We will build a temporary chat session using the history
    messages = []
    for msg in chat_history[:-1]:
        messages.append({
            "role": "user" if msg["role"] == "user" else "model",
            "parts": [msg["content"]]
        })
    
    messages.append({
        "role": "user",
        "parts": [f"{context_prefix}我的問題是：{query}"]
    })
    
    try:
        response = _generate_with_fallback(
            api_key=api_key,
            system_instruction=system_instruction,
            messages=messages
        )
        return response.text
    except Exception as e:
        # Try direct text generation without chat wrapper as last resort
        try:
            direct_prompt = f"{system_instruction}\n\n歷史對話內容:\n"
            for msg in chat_history[:-1]:
                role_label = "學員" if msg["role"] == "user" else "導師"
                direct_prompt += f"{role_label}: {msg['content']}\n"
            direct_prompt += f"\n目前問題:\n{context_prefix}問題：{query}"
            
            response = _generate_with_fallback(
                api_key=api_key,
                prompt=direct_prompt
            )
            return response.text
        except Exception as e2:
            raise RuntimeError(f"AI Tutor Chat failed: {e2}")

def generate_rpg_scenario(filename, api_key):
    import database
    import random
    chunks = database.get_chunks_for_file(filename)
    if not chunks:
        raise ValueError("此教材尚未被向量化或無內容。")
    random_chunk = random.choice(chunks)[0]
    
    rules = database.load_ai_rules()
    rules_text = "\n".join([f"- {r}" for r in rules]) if rules else "無"
    
    system_instruction = (
        "你是一位嚴格但專業的麻醉科/急診科主治醫師 (Dungeon Master)。"
        "請根據提供的【教材內容】，設計一個寫實的突發臨床情境。"
        "情境必須包含病患的基本資訊（年齡、性別）與突發狀況的表徵。"
        "結尾必須詢問護理生：『請問你現在要採取什麼行動？』"
        "絕對不要在情境中透露解答。\n\n"
        "【來自總醫師的絕對指令與糾錯】：\n"
        f"{rules_text}\n"
    )
    prompt = f"【教材內容】\n{random_chunk}\n\n請產生一個臨床情境開場白。"
    
    response = _generate_with_fallback(
        api_key=api_key,
        system_instruction=system_instruction,
        prompt=prompt
    )
    return response.text

def evaluate_rpg_action(filename, chat_history, api_key):
    query_text = " ".join([m["content"] for m in chat_history[-3:]])
    context_chunks = search_related_chunks(query_text, filename, api_key=api_key, top_k=3)
    context_text = "\n\n".join([c["text"] for c in context_chunks])
    
    import database
    rules = database.load_ai_rules()
    rules_text = "\n".join([f"- {r}" for r in rules]) if rules else "無"
    
    system_instruction = (
        "你是一位溫和、具備同理心且熱心教學的麻醉科/急診科主治醫師 (DM)，現在是『友善引導模式』。"
        "根據學生的處置與【標準教材內容】，給予正向的結果與回饋。"
        "如果學生的處置大方向正確，請給予肯定，並明確在結尾回覆【模擬結束】四個字。"
        "如果處置不完整或錯誤，請溫柔地提醒他可能忽略的地方，病患狀況會稍微改變但不會立即死亡，並再次詢問『你下一步要怎麼做？』"
        "絕對不要直接給解答，必須引導學生自己思考。"
        "必須遵守標準教材內容，不要自行發明處置流程。\n\n"
        "【來自總醫師的絕對指令與糾錯】：\n"
        f"{rules_text}\n"
    )
    
    messages = []
    for msg in chat_history[:-1]:
        messages.append({
            "role": "user" if msg["role"] == "user" else "model",
            "parts": [msg["content"]]
        })
        
    last_msg = chat_history[-1]["content"]
    messages.append({
        "role": "user",
        "parts": [f"【標準教材內容】\n{context_text}\n\n學生的處置：{last_msg}"]
    })
    
    response = _generate_with_fallback(
        api_key=api_key,
        system_instruction=system_instruction,
        messages=messages
    )
    
    res_text = response.text
    is_finished = "模擬結束" in res_text
    return res_text, is_finished
